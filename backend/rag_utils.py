"""
rag_utils.py
Everything needed to turn an uploaded document into searchable vector
chunks, and to turn a user's question into a ranked list of relevant
chunks + a generated answer.
"""

from __future__ import annotations

import os
from typing import List, Optional, Sequence
from uuid import UUID

from sqlalchemy.orm import Session

import models

# Lazy-loaded so importing routers does not block on API client init.
_gemini_client = None
_groq_client = None

# Keep this in sync with the `embedding` column's Vector(...) dimension in models.py.
EMBED_DIM = 768


def check_rag_dependencies() -> tuple[bool, str]:
    """
    Verify packages needed to index documents and answer from them.
    Uses find_spec so startup does not load heavy libraries.
    """
    import importlib.util

    missing = []
    for mod_name, pip_name in (
        ("google.genai", "google-genai"),
        ("pypdf", "pypdf"),
        ("docx", "python-docx"),
        ("pptx", "python-pptx"),
        ("groq", "groq"),
    ):
        if importlib.util.find_spec(mod_name) is None:
            missing.append(pip_name)
    if missing:
        return False, (
            "RAG packages missing: "
            + ", ".join(missing)
            + ". Install with: pip install -r requirements-rag.txt"
        )
    if not os.getenv("GEMINI_API_KEY"):
        return False, (
            "GEMINI_API_KEY is not set — document indexing will fail "
            "(embeddings use the Gemini API)."
        )
    if not os.getenv("GROQ_API_KEY"):
        return False, (
            "GROQ_API_KEY is not set — document indexing works, "
            "but chat answers will fail."
        )
    return True, "RAG ready"


def _get_gemini_client():
    global _gemini_client
    if _gemini_client is None:
        try:
            from google import genai
        except ImportError as exc:
            raise RuntimeError(
                "google-genai is not installed. "
                "Run: pip install -r requirements-rag.txt"
            ) from exc

        _gemini_client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))
    return _gemini_client


def _get_groq_client():
    global _groq_client
    if _groq_client is None:
        try:
            from groq import Groq
        except ImportError as exc:
            raise RuntimeError(
                "groq is not installed. Run: pip install -r requirements-rag.txt"
            ) from exc

        _groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))
    return _groq_client


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def extract_text(file_path: str) -> str:
    if not file_path or not os.path.isfile(file_path):
        return ""

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        text = _extract_pdf(file_path)
    elif ext == ".docx":
        text = _extract_docx(file_path)
    elif ext == ".pptx":
        text = _extract_pptx(file_path)
    elif ext in (".txt", ".csv", ".md", ".log"):
        text = _extract_plain(file_path)
    else:
        # Unsupported type — skip RAG indexing rather than fail the upload.
        return ""

    return _normalize_extracted_text(text)


def _normalize_extracted_text(text: str) -> str:
    """Collapse runaway blank lines so embeddings are not dominated by whitespace."""
    if not text:
        return ""
    lines = [ln.rstrip() for ln in text.splitlines()]
    compact: list[str] = []
    blank_run = 0
    for ln in lines:
        if not ln.strip():
            blank_run += 1
            if blank_run <= 1:
                compact.append("")
            continue
        blank_run = 0
        compact.append(ln)
    return "\n".join(compact).strip()


def _extract_pdf(file_path: str) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError(
            "pypdf is not installed. Run: pip install pypdf"
        ) from exc
    reader = PdfReader(file_path)
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_docx(file_path: str) -> str:
    try:
        from docx import Document as DocxDocument
    except ImportError as exc:
        raise RuntimeError(
            "python-docx is not installed. Run: pip install python-docx"
        ) from exc
    doc = DocxDocument(file_path)
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_plain(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def _extract_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is not installed. Run: pip install python-pptx"
        ) from exc

    prs = Presentation(file_path)
    lines: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    # Prefer para.text — some PPTX files store content outside runs.
                    text = (para.text or "").strip()
                    if not text:
                        text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    # Skip empty rows — they used to inject blank lines that
                    # wrecked chunk embeddings (title slides especially).
                    if any(cells):
                        slide_lines.append(" | ".join(cells))
        if slide_lines:
            lines.append(f"--- Slide {idx} ---")
            lines.extend(slide_lines)
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------

def chunk_text(text: str, chunk_size: int = 800, overlap: int = 150) -> List[str]:
    """
    Simple sliding-window character chunker. Good enough for most documents;
    swap for a token-aware / sentence-aware splitter later if chunk
    boundaries cut sentences awkwardly in practice.
    """
    text = text.strip()
    if not text:
        return []

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += chunk_size - overlap

    return [c.strip() for c in chunks if c.strip()]


# ---------------------------------------------------------------------------
# Embeddings (Gemini API — no local model download)
# ---------------------------------------------------------------------------

def get_embedding(text: str, *, task_type: str = "RETRIEVAL_DOCUMENT") -> list[float]:
    """
    task_type should be "RETRIEVAL_DOCUMENT" when indexing chunks and
    "RETRIEVAL_QUERY" when embedding a user's question — Gemini tunes the
    embedding differently for each, which improves retrieval quality.
    """
    from google.genai.types import EmbedContentConfig

    response = _get_gemini_client().models.embed_content(
        model="gemini-embedding-001",
        contents=text,
        config=EmbedContentConfig(
            task_type=task_type,
            output_dimensionality=EMBED_DIM,
        ),
    )
    return response.embeddings[0].values


def chunk_count_for_document(db: Session, document_id) -> int:
    return (
        db.query(models.DocumentChunk)
        .filter(models.DocumentChunk.document_id == document_id)
        .count()
    )


# ---------------------------------------------------------------------------
# Ingestion — call this after a document is saved to disk
# ---------------------------------------------------------------------------

def process_document_for_rag(
    db: Session,
    document_id,
    file_path: str,
    organization_id,
    project_id=None,
) -> int:
    """
    Returns the number of chunks created (0 if the file type is unsupported
    or the file has no extractable text).

    Safe to call from a request handler (upload / reindex). Also available via
    `process_document_for_rag_background` when a separate DB session is needed.
    """
    try:
        text = extract_text(file_path)
    except Exception as exc:
        print(f"RAG extraction failed for document {document_id}: {exc}")
        return 0

    if not text or not text.strip():
        print(f"RAG: no extractable text for document {document_id} ({file_path})")
        return 0

    # Replace any previous chunks for this document (re-index / retry safe).
    db.query(models.DocumentChunk).filter(
        models.DocumentChunk.document_id == document_id
    ).delete(synchronize_session=False)

    pieces = chunk_text(text)
    for piece in pieces:
        embedding = get_embedding(piece, task_type="RETRIEVAL_DOCUMENT")
        db.add(models.DocumentChunk(
            document_id=document_id,
            organization_id=organization_id,
            project_id=project_id,
            chunk_text=piece,
            embedding=embedding,
        ))
    db.commit()
    print(f"RAG: indexed {len(pieces)} chunk(s) for document {document_id}")
    return len(pieces)


def process_document_for_rag_background(
    document_id,
    file_path: str,
    organization_id,
    project_id=None,
) -> int:
    """
    Background-task entrypoint. Opens and closes its own SessionLocal so it
    does not depend on the request-scoped session (which FastAPI closes after
    the upload response is sent).
    """
    from database import SessionLocal

    db = SessionLocal()
    try:
        return process_document_for_rag(
            db,
            document_id,
            file_path,
            organization_id,
            project_id,
        )
    except Exception as exc:
        print(
            f"RAG background indexing failed for document {document_id}: "
            f"{type(exc).__name__}: {exc}"
        )
        try:
            db.rollback()
        except Exception:
            pass
        return 0
    finally:
        db.close()


def delete_chunks_for_document(db: Session, document_id, *, commit: bool = True):
    """Call this whenever a document is deleted, so stale chunks don't
    linger and get retrieved in future answers."""
    db.query(models.DocumentChunk).filter(
        models.DocumentChunk.document_id == document_id
    ).delete(synchronize_session=False)
    if commit:
        db.commit()


# ---------------------------------------------------------------------------
# Retrieval
# ---------------------------------------------------------------------------

_RETRIEVAL_STOPWORDS = frozenset(
    "a an the is are was were be been being what which who whom this that "
    "these those am do does did doing have has had having of in on to for "
    "with about into through during before after above below from up down "
    "out off over under again further then once here there when where why "
    "how all any both each few more most other some such no nor not only "
    "own same so than too very can will just should now i me my we our you "
    "your he she it they them their".split()
)


def _keyword_overlap(question: str, chunk_text: str) -> float:
    terms = {
        t
        for t in "".join(
            ch.lower() if ch.isalnum() else " " for ch in question
        ).split()
        if len(t) > 2 and t not in _RETRIEVAL_STOPWORDS
    }
    if not terms:
        return 0.0
    hay = chunk_text.lower()
    hits = sum(1 for t in terms if t in hay)
    return hits / len(terms)


def _hybrid_score(
    query_embedding: Sequence[float],
    chunk,
    question: str,
    *,
    earliest_ids: set,
) -> float:
    emb = chunk.embedding
    if not emb:
        vec = -1.0
    else:
        try:
            vec = _cosine_similarity(query_embedding, list(emb))
        except Exception:
            vec = -1.0
    kw = _keyword_overlap(question, chunk.chunk_text or "")
    # First chunk of a document is often the title / abstract page.
    early = 0.08 if chunk.id in earliest_ids else 0.0
    return float(vec) + 0.22 * kw + early


def retrieve_relevant_chunks(
    db: Session,
    question: str,
    organization_id,
    project_id=None,
    document_ids: Optional[Sequence[UUID]] = None,
    k: int = 5,
):
    """
    Rank document chunks by cosine distance to the question embedding,
    then lightly re-rank with keyword overlap so title/intro pages are not
    missed on short factual questions.

    Scoping rules:
    - Always limited to the caller's organization.
    - If `document_ids` is provided (even empty), only those documents are
      searched. An empty list means "no documents selected" → no results.
    - Else if `project_id` is set, only that project's chunks are searched.
    - Else all org chunks are searched.
    """
    query_embedding = get_embedding(question, task_type="RETRIEVAL_QUERY")

    query = db.query(models.DocumentChunk).filter(
        models.DocumentChunk.organization_id == organization_id
    )

    if document_ids is not None:
        if not document_ids:
            return []
        query = query.filter(models.DocumentChunk.document_id.in_(list(document_ids)))
    elif project_id:
        query = query.filter(models.DocumentChunk.project_id == project_id)

    candidate_limit = max(k * 6, 24)
    candidates = []

    # Prefer vector ranking when pgvector is available; otherwise rank in
    # Python (JSON embeddings / SQLite / extension missing).
    try:
        candidates = (
            query.order_by(
                models.DocumentChunk.embedding.cosine_distance(query_embedding)
            )
            .limit(candidate_limit)
            .all()
        )
    except Exception as exc:
        print(f"RAG vector search unavailable ({exc}); ranking in Python.")
        # Failed SQL leaves the transaction aborted on Postgres — reset before
        # the Python fallback can query again.
        try:
            db.rollback()
        except Exception:
            pass
        query = db.query(models.DocumentChunk).filter(
            models.DocumentChunk.organization_id == organization_id
        )
        if document_ids is not None:
            query = query.filter(
                models.DocumentChunk.document_id.in_(list(document_ids))
            )
        elif project_id:
            query = query.filter(models.DocumentChunk.project_id == project_id)
        candidates = query.limit(max(k * 40, 200)).all()

    if not candidates:
        return []

    # Pin the earliest chunk per document into the candidate pool (title pages).
    by_doc: dict = {}
    for chunk in candidates:
        by_doc.setdefault(chunk.document_id, []).append(chunk)

    earliest_ids: set = set()
    pinned = []
    for doc_id in by_doc:
        first = (
            db.query(models.DocumentChunk)
            .filter(models.DocumentChunk.document_id == doc_id)
            .order_by(models.DocumentChunk.created_at.asc())
            .first()
        )
        if first is not None:
            earliest_ids.add(first.id)
            if all(c.id != first.id for c in candidates):
                pinned.append(first)
    if pinned:
        candidates = list(candidates) + pinned

    scored = [
        (_hybrid_score(query_embedding, chunk, question, earliest_ids=earliest_ids), chunk)
        for chunk in candidates
    ]
    scored.sort(key=lambda pair: pair[0], reverse=True)

    seen = set()
    out = []
    for _, chunk in scored:
        if chunk.id in seen:
            continue
        seen.add(chunk.id)
        out.append(chunk)
        if len(out) >= k:
            break
    return out


def _cosine_similarity(a: Sequence[float], b: Sequence[float]) -> float:
    if not a or not b or len(a) != len(b):
        return -1.0
    dot = sum(x * y for x, y in zip(a, b))
    return float(dot)  # embeddings are L2-normalized in get_embedding()


# ---------------------------------------------------------------------------
# Generation — Groq is used for the final answer. Nothing else in the RAG
# pipeline needs to change since everything else calls call_llm(), not Groq
# directly.
# ---------------------------------------------------------------------------

def call_llm(prompt: str) -> str:
    if not os.getenv("GROQ_API_KEY"):
        return (
            "Sorry, I couldn't generate an answer right now "
            "(GROQ_API_KEY is not configured)."
        )
    try:
        response = _get_groq_client().chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=500,
            temperature=0.2,
        )
        return response.choices[0].message.content
    except Exception as exc:
        return f"Sorry, I couldn't generate an answer right now ({exc.__class__.__name__})."


# ---------------------------------------------------------------------------
# End-to-end RAG answer — used by chat_assistant_service as the fallback when
# no keyword intent matches.
# ---------------------------------------------------------------------------

def warmup_embedding_model() -> None:
    """Make one cheap embedding call at startup so the first real chat isn't
    the one paying for client init / first-call latency."""
    try:
        get_embedding("warmup", task_type="RETRIEVAL_QUERY")
        print("RAG embedding client warmed up.")
    except Exception as exc:  # pragma: no cover
        print(f"RAG embedding warmup skipped: {exc}")


def get_rag_response(
    db: Session,
    question: str,
    organization_id,
    project_id: Optional[UUID] = None,
    document_ids: Optional[Sequence[UUID]] = None,
    k: int = 5,
) -> str:
    # Normalize string UUIDs from JSON clients.
    normalized_ids: Optional[list[UUID]] = None
    if document_ids is not None:
        normalized_ids = []
        for raw in document_ids:
            if raw is None:
                continue
            normalized_ids.append(raw if isinstance(raw, UUID) else UUID(str(raw)))

    if normalized_ids is not None and len(normalized_ids) > 0:
        indexed = (
            db.query(models.DocumentChunk.document_id)
            .filter(models.DocumentChunk.document_id.in_(normalized_ids))
            .distinct()
            .all()
        )
        indexed_ids = {row[0] for row in indexed}
        missing = [doc_id for doc_id in normalized_ids if doc_id not in indexed_ids]
        if not indexed_ids:
            names = (
                db.query(models.Document.filename)
                .filter(models.Document.id.in_(normalized_ids))
                .all()
            )
            label = ", ".join(n[0] for n in names) or "selected file(s)"
            return (
                f"Those documents are not indexed for chat yet ({label}). "
                "Open Documents → click Reindex on each file, wait until it "
                "shows chunk count > 0, then ask again."
            )
        # Search only indexed selections; ignore unindexed duplicates.
        if missing:
            normalized_ids = [doc_id for doc_id in normalized_ids if doc_id in indexed_ids]

    chunks = retrieve_relevant_chunks(
        db,
        question,
        organization_id,
        project_id=project_id,
        document_ids=normalized_ids,
        k=k,
    )
    if not chunks:
        if normalized_ids is not None and len(normalized_ids) == 0:
            return (
                "No documents are available to search. Upload a PDF, DOCX, "
                "PPTX, or text file on the Documents page, then ask again."
            )
        if normalized_ids is not None:
            return (
                "I couldn't find anything relevant in those documents. "
                "Try rephrasing the question, or Reindex on the Documents page "
                "if the file shows as not indexed. Supported types: "
                "PDF, DOCX, PPTX, TXT, CSV, MD."
            )
        return (
            "I couldn't find anything relevant in your documents. "
            "Upload a PDF, DOCX, PPTX, or text file first, then ask again."
        )

    # Attach filenames so the model (and user) can see sources.
    doc_ids = {c.document_id for c in chunks}
    filename_by_id = {}
    if doc_ids:
        docs = (
            db.query(models.Document)
            .filter(models.Document.id.in_(list(doc_ids)))
            .all()
        )
        filename_by_id = {d.id: d.filename for d in docs}

    context_parts = []
    sources = []
    for c in chunks:
        name = filename_by_id.get(c.document_id, "unknown file")
        context_parts.append(f"[Source: {name}]\n{c.chunk_text}")
        if name not in sources:
            sources.append(name)

    context = "\n\n---\n\n".join(context_parts)
    prompt = f"""Answer the question using only the context below.
If the answer isn't in the context, say you don't have that information.
Be concise and cite the source filename when helpful.

Context:
{context}

Question: {question}"""

    answer = call_llm(prompt)
    if sources:
        answer = f"{answer}\n\nSources: {', '.join(sources)}"
    return answer